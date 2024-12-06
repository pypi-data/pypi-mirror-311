import os
import platform
import re
import shutil
import subprocess
import zipfile
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime
from glob import glob

import extract_msg
import pandas as pd
import pptx
from docx import Document
from pypdf import PdfReader
from tqdm import tqdm

"""

This file is used to store code for extracting text from different file formats.

@author: Michael de Winter
"""


def is_office_installed():
    """
    Check if there is a office installation.
    """
    # Common installation directories for Microsoft Office
    common_paths = [
        r"C:\Program Files\Microsoft Office",
        r"C:\Program Files (x86)\Microsoft Office",
        r"C:\Program Files\Microsoft Office\root\Office16",  # Office 2016/365
        r"C:\Program Files (x86)\Microsoft Office\root\Office16",
    ]

    for path in common_paths:
        if os.path.exists(path):
            print(f"Microsoft Office is installed at {path}.")
            return True

    print("Microsoft Office is not installed.")
    return False


def is_libre_office_installed():
    """
    Check if there is LibreOffice installation.

    """
    return os.path.isfile("C:/Program Files/LibreOffice/program/soffice.exe")


def process_and_save_file(
    afilepath,
    atarget_dir,
    atarget_dir_extracted,
    verbose=False,
):
    """

    Method to call process_file and correctly generate a filename if it contains subdirectories.

    @param afilepath: Path to a specific file.
    @param atarget_dir: the original directory of the file is stripped from the filename for naming purposes to get it's subdirectory.
    @oaram atarget_dir_extracted: directory in which the file will be stored.
    """

    if verbose:
        print("-------")
        print(afilepath)
    afilepath = afilepath.replace("\\", "/")

    # remove linebreaks with "-"
    text_content = re.sub(r"-\s*\n\s*", "", process_file(afilepath))
    # Normalize a file by removing \b and whitespaces
    text_content = re.sub(r"\s+", " ", text_content)

    try:
        # Generate a output name but if there are subdirectories in the directory put the subdirectories name in the filename.
        a_output_name = (
            atarget_dir_extracted
            + "/"
            + re.sub(
                r"\.(?!.*\.)",
                "_",
                afilepath.replace(atarget_dir, "").replace("/", "#"),
            )
            + ".txt"
        )
        if verbose:
            print(a_output_name)
        with open(a_output_name, "w", encoding="utf-8") as file:
            file.write(str(text_content))
            file.close()
    except Exception as e:
        print(f"Error with writing file {afilepath}: " + e)


def make_dir_from_filename(afilepath):
    """

    Make a directory based on the file name

    """
    a_output_directory = afilepath.rsplit(".", 1)[0]
    print("Making :" + str(a_output_directory))
    os.makedirs(a_output_directory)
    return a_output_directory


def extract_msg_attachments(atarget_dir):
    """
    Extract all .msg i.e. emails files in a directory.

    @param atarget_dir: directory where the .msg files are stored.
    """
    # Extract all .msg files into a directory
    for afilepath in glob(atarget_dir + "*.msg"):
        afilepath = afilepath.replace("\\", "/")

        encodings = ["utf-8", "ISO-8859-1", "windows-1252"]
        for encoding in encodings:
            try:
                msg = extract_msg.Message(afilepath, overrideEncoding=encoding)

                if len(msg.attachments) > 0:
                    a_output_directory = make_dir_from_filename(afilepath)

                    # TODO: Check if encoding is also needed for the attachments.
                    for item in range(0, len(msg.attachments)):
                        att = msg.attachments[item]
                        msg.attachments[item].save(
                            customPath=a_output_directory,
                            customFilename=att.longFilename,
                        )
                return True
            except Exception as e:
                continue
        print(f"Failed to read {afilepath} with all tested encodings.")
        return False


def extract_zip_attachments(atarget_dir):
    """

    Extract all .zip files in a directory.

    @param atarget_dir: directory where the .zip files are stored.
    """
    # Extract a zip file into a new directory

    for afilepath in glob(atarget_dir + "*.zip"):
        try:
            afilepath = afilepath.replace("\\", "/")
            print(afilepath)
            a_output_directory = make_dir_from_filename(afilepath)

            # Extract the contents of the zip file
            with zipfile.ZipFile(afilepath, "r") as zip_ref:
                zip_ref.extractall(a_output_directory)
        except Exception as e:
            print(e)


def extracted_files_from_list_filepaths(
    afilepaths, atarget_dir, atarget_dir_extracted, verbose=False, threads=False
):
    """

    Method to extract all files in directory in a .txt format in a other directory for easy searching and snapshotting.

    @param afilepaths: a array of filepaths to be extracted to .txt files.
    @param atarget_dir: the original directory of the file is stripped from the filename for naming purposes to get it's subdirectory.
    @oaram atarget_dir_extracted: directory in which the file will be stored.
    @param verbose: Whether to print more output.
    @param threads: Whether to multiprocess with the file reading thus speeding up the process but not all systems can handle it.
    """
    if threads:
        with ThreadPoolExecutor() as executor:
            futures = {
                executor.submit(
                    process_and_save_file,
                    afilepath,
                    atarget_dir,
                    atarget_dir_extracted,
                    verbose,
                ): afilepath
                for afilepath in afilepaths
            }
            for future in tqdm(as_completed(futures), total=len(futures)):
                try:
                    future.result()
                except Exception as e:
                    print(f"Exception occurred: {e}")
    else:
        for afilepath in tqdm(afilepaths):
            process_and_save_file(
                afilepath, atarget_dir, atarget_dir_extracted, verbose
            )


def run(atarget_dir, atarget_dir_extracted, amulti_thread=False):
    """
    Method to run the process.

    @param atarget_dir: the original directory of the file is stripped from the filename for naming purposes to get it's subdirectory.
    @oaram atarget_dir_extracted: directory in which the file will be stored.
    @param amulti_thread: Whether to multiprocess with the file reading thus speeding up the process but not all systems can handle it. TODO: Find out why
    """
    # First extract all the .msg files.
    print("Extracting .msg files")
    extract_msg_attachments(atarget_dir)
    # Then extract all the .zip files.
    print("Extracting .zip files")
    extract_zip_attachments(atarget_dir)

    # Make a directory if the directory does not exist yet.
    if os.path.isdir(atarget_dir_extracted) == False:
        os.makedirs(atarget_dir_extracted)

    extracted_files_from_list_filepaths(
        [
            afilepath
            for afilepath in glob(os.path.join(atarget_dir, "**", "*"), recursive=True)
        ],
        atarget_dir,
        atarget_dir_extracted,
        threads=amulti_thread,
    )


# Different functions for opening files.
def read_pdf(file_path, detect_encoding=False):
    text = ""
    try:
        with open(file_path, "rb") as file:
            # Create a PDF reader object
            pdf_reader = PdfReader(file)

            # Get the number of pages in the PDF
            num_pages = len(pdf_reader.pages)

            # Extract text from each page
            for page_num in range(num_pages):
                page = pdf_reader.pages[page_num]

                # Raw text extraction
                raw_text = page.extract_text()

                # Analyze encoding
                encoded_text = raw_text

                if detect_encoding == True:
                    # TODO: Make this run with chardet.
                    # detected_encoding = chardet.detect(encoded_text)  # Detect encoding

                    # Log detected encoding
                    # print(
                    #    f"Page {page_num + 1}: Detected encoding - {detected_encoding['encoding']}"
                    # )

                    # Decode and add to full text if encoding is valid
                    # try:
                    #    decoded_text = encoded_text.decode(
                    #        detected_encoding["encoding"], errors="ignore"
                    #    )
                    #    text += decoded_text
                    # except Exception as e:
                    #    print(f"Error decoding page {page_num + 1}: {e}")
                    print("Error not yet implementated")
                else:
                    text += encoded_text

    except Exception as e:
        print(f"Error processing {file_path}: {e}")

    return text


def read_docx(file_path):
    doc = Document(file_path)
    text = []
    for paragraph in doc.paragraphs:
        text.append(paragraph.text)
    return "\n".join(text)


def read_msg(file_path):
    encodings = ["utf-8", "ISO-8859-1", "windows-1252"]
    for encoding in encodings:
        try:
            msg = extract_msg.Message(file_path, overrideEncoding=encoding)
            msg_text = "Sender: " + str(msg.sender) + " | \n "
            msg_text += "To: " + str(msg.to) + " | \n "
            msg_text += "CC: " + str(msg.cc) + " | \n "
            msg_text += "BCC: " + str(msg.bcc) + " | \n "
            msg_text += "Subject: " + str(msg.subject) + " | \n "
            msg_text += "Body: " + str(msg.body)
            return msg_text  # Return if successful
        except Exception as e:
            continue  # Try the next encoding if an error occurs

    print(f"Failed to read {file_path} with all tested encodings.")
    return None  # Return None if all encoding attempts fail


def read_xls(file_path):
    # Convert the pandas DataFrame to a single string
    data = pd.read_excel(file_path)
    collapsed_text = data.to_string(index=False)
    return collapsed_text


def detect_os():
    """Detect the operating system."""
    system = platform.system()
    if system == "Windows":
        return "Windows"
    elif system == "Linux" or system == "Darwin":
        return "Linux/Unix"
    else:
        raise OSError("Unsupported operating system")


def process_doc_with_libreoffice(file_path):
    """
    Extract text from a .doc file using LibreOffice and remove the generated .txt file.
    ENSURE THAT LibreOffice is installed.

    @param file_path (str): Path to the .doc file.

    Returns:
        str: Extracted text from the .doc file.
    """
    try:
        # Define output directory and file name

        output_file = file_path.replace(".doc", "")

        # Set path to libre office.
        subprocess.run(
            "set PATH=%PATH%;C:/Program Files/LibreOffice/program", shell=True
        )

        # Convert .doc to .txt using LibreOffice
        subprocess.run(
            [
                "soffice",
                "--headless",
                "--convert-to",
                "txt",
                "--outdir",
                output_file,
                file_path,
            ],
            check=True,
        )

        # Read and store the content of the generated .txt file
        with open(glob(output_file + "/*")[0], "r", encoding="utf-8") as txt_file:
            content = txt_file.read()

        # Remove the generated .txt file
        shutil.rmtree(output_file)

        return content
    except subprocess.CalledProcessError as e:
        print(f"LibreOffice conversion error: {e}")
        return None
    except FileNotFoundError:
        print("LibreOffice is not installed or not in PATH.")
        return None
    except Exception as e:
        print(f"An error occurred: {e}")
        return None


def process_doc_with_pywin32(file_path):
    """Process .doc file using pywin32 on Windows."""
    try:
        # Import win32com only if on Windows
        import win32com.client

        # Initialize Word application
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False  # Prevent Word GUI

        # Open the .doc file
        doc = word.Documents.Open(os.path.abspath(file_path))

        # Extract text
        text = doc.Content.Text

        # Close the document and Word application
        doc.Close()
        word.Quit()

        return text

    except Exception as e:
        raise RuntimeError(f"Error processing .doc file with pywin32: {e}")


def process_doc_with_antiword(file_path):
    """
    Extract text from a .doc file using Antiword.
    :param file_path: Path to the .doc file.
    :return: Extracted text.
    """
    try:
        # Call Antiword to extract text
        result = subprocess.run(
            ["antiword", file_path],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
        )

        # Check for errors
        if result.returncode != 0:
            raise RuntimeError(f"Antiword error: {result.stderr.strip()}")

        # Return the extracted text
        return result.stdout

    except FileNotFoundError:
        raise RuntimeError("Antiword is not installed. Please install it first.")
    except Exception as e:
        raise RuntimeError(f"An error occurred while processing the file: {e}")


def read_doc(file_path):
    """Read .doc file based on the operating system."""
    os_type = detect_os()

    if os_type == "Windows":
        print("Detected OS: Windows. Using pywin32 to process the .doc file.")

        # TODO:Might optimize performance here a bit with not constantly checking install locations.
        # Detect if libreOffice is installed at the default location.
        if is_libre_office_installed:
            return process_doc_with_libreoffice(file_path)
        elif is_office_installed:
            return process_doc_with_pywin32(file_path)
        else:
            raise "Windows Office or LibreOffice not found on the default locations to process .doc"
    elif os_type == "Linux/Unix":
        print("Detected OS: Linux/Unix. Using unoconv to process the .doc file.")
        return process_doc_with_antiword(file_path)
    else:
        raise OSError("Unsupported operating system for .doc processing")


def read_pptx(file_path):
    text = ""
    presentation = pptx.Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text
    return text


def process_file(a_file_path):
    """

    Open a file based on a file extension.

    @param a_file_path: File path to a file.
    """

    file_split = a_file_path.split(".")
    file_type = a_file_path.rsplit(".", 1)[-1]
    a_content = ""

    if len(file_split) > 1:
        try:
            if "pdf" in file_type:
                # Action for PDF file type
                a_content = read_pdf(a_file_path)
            elif "xls" in file_type or "xlsx" in file_type:
                # Action for XLS file type
                a_content = read_xls(a_file_path)
            elif "docx" in file_type:
                # Action for DOCX file type
                a_content = read_docx(a_file_path)
            elif "doc" in file_type:
                a_content = read_doc(a_file_path)
            elif "msg" in file_type:
                # Action for MSG file type
                a_content = read_msg(a_file_path)
            elif "txt" in file_type:
                with open(a_file_path, "r") as file:
                    # Read the content of the file
                    a_content = file.read()
            elif "pptx" in file_type:
                a_content = read_pptx(a_file_path)
        except Exception as e:
            print("Error with Reading " + a_file_path)
            print(e)

    return a_content


def generate_filename(output_dir):
    """
    Generate a filename output based on the current date

    @param output_dir: The directory to base the filename on.
    """
    current_date = datetime.now().strftime("%Y%m%d_%H%M%S")  # Format as %Y%m%d_%H%M%S
    filename = f"{output_dir}/selected_word_counter_output_{current_date}"
    return filename


def replace_last_slash(path, replacement=""):
    """
    This function is used for testing.

    @param path: path to be split.
    @param replacement: value to replace / with.
    """
    parts = path.rsplit("/", 1)
    return replacement.join(parts)


def delete_directory(path):
    try:
        shutil.rmtree(path)
        print(f"Directory '{path}' and all its contents were deleted successfully.")
    except FileNotFoundError:
        print(f"Directory '{path}' not found.")
    except Exception as e:
        print(f"An error occurred: {e}")
