import logging
from math import floor
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Mm, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER


# Utility Functions
def _set_text_properties(text_frame, text, font_props, alignment=PP_ALIGN.LEFT):
    """
    Sets the text and font properties for a given text frame.
    """
    text_frame.clear()  # Clear existing text
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.name = font_props.get('name', 'Arial')
    font.size = font_props.get('size', Pt(14))
    font.bold = font_props.get('bold', False)
    font.color.rgb = font_props.get('color', RGBColor(0, 0, 0))  # Default to black
    paragraph.alignment = alignment  # Set alignment


def _get_placeholder_by_type(slide, placeholder_type):
    """
    Retrieves a placeholder from the slide by its placeholder type.
    """
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        if shape.placeholder_format.type == placeholder_type:
            return shape
    return None


def _add_table(slide, dataframe, left_mm, top_mm, width_mm, height_mm, header_font_props, cell_font_props):
    """
    Adds a table to the slide based on a pandas DataFrame.
    """
    rows, cols = dataframe.shape
    left = Cm(left_mm / 10)
    top = Cm(top_mm / 10)
    width = Cm(width_mm / 10)
    height = Cm(height_mm / 10)
    table_shape = slide.shapes.add_table(rows + 1, cols, left, top, width, height)
    table = table_shape.table

    # Apply header
    for col_idx, column_name in enumerate(dataframe.columns):
        cell = table.cell(0, col_idx)
        _set_text_properties(cell.text_frame, str(column_name), header_font_props, alignment=PP_ALIGN.CENTER)

    # Apply cell data
    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx + 1, col_idx)
            _set_text_properties(cell.text_frame, str(dataframe.iloc[row_idx, col_idx]), cell_font_props,
                                 alignment=PP_ALIGN.CENTER)


class CustomPresentationBuilder:
    def __init__(self, slides_content_list):
        self.slides_content_list = slides_content_list
        self.prs = Presentation()
        # Set slide size to 16:9 aspect ratio in centimeters
        self.prs.slide_width = Cm(33.867)  # 13.3333 inches * 2.54 cm/in ≈ 33.867 cm
        self.prs.slide_height = Cm(19.05)  # 7.5 inches * 2.54 cm/in = 19.05 cm
        self.layout_mapping = self._create_layout_mapping()
        self.default_fonts = {
            'title': {'name': 'Arial', 'size': Pt(26.7), 'bold': False, 'color': RGBColor(0, 0, 0)},
            'body': {'name': 'Arial', 'size': Pt(14), 'bold': False, 'color': RGBColor(0, 0, 0)},
            'header': {'name': 'ＭＳ Ｐゴシック (本文)', 'size': Pt(11), 'bold': True, 'color': RGBColor(0, 0, 0)},
            'cell': {'name': 'ＭＳ Ｐゴシック (本文)', 'size': Pt(11), 'bold': False, 'color': RGBColor(0, 0, 0)},
        }

    def _create_layout_mapping(self):
        layout_mapping = {}
        for layout in self.prs.slide_layouts:
            name = layout.name.upper().replace(' ', '_')
            layout_mapping[name] = layout
        return layout_mapping

    def create_presentation(self, output_file):
        for idx, slide_content in enumerate(self.slides_content_list, start=1):
            layout_name = slide_content.get('layout', 'TITLE_AND_CONTENT').upper()
            layout = self.layout_mapping.get(layout_name, self.prs.slide_layouts[1])  # Default to 'TITLE_AND_CONTENT'
            slide = self.prs.slides.add_slide(layout)
            print(f"Creating Slide {idx}: Using layout '{layout.name}'")

            # Handle Title
            title_text = slide_content.get('title', '')
            title_font_props = slide_content.get('title_font', self.default_fonts['title'])
            title_position = slide_content.get('title_position', None)  # Optional

            if title_position:
                # Add title as a textbox with specified position
                left = Cm(title_position.get('left_cm', 2))
                top = Cm(title_position.get('top_cm', 1))
                width = Cm(title_position.get('width_cm', 29.867))  # Default slide width minus margins
                height = Cm(title_position.get('height_cm', 3))
                txBox = slide.shapes.add_textbox(left, top, width, height)
                _set_text_properties(txBox.text_frame, title_text, title_font_props,
                                     alignment=title_position.get('alignment', PP_ALIGN.LEFT))
                print(f"Added custom title textbox at ({left.cm} cm, {top.cm} cm)")
            else:
                # Use existing title placeholder if available
                title_placeholder = _get_placeholder_by_type(slide, PP_PLACEHOLDER.TITLE)
                if not title_placeholder:
                    title_placeholder = slide.shapes.title  # Fallback
                if title_placeholder and title_placeholder.text_frame:
                    _set_text_properties(title_placeholder.text_frame, title_text, title_font_props,
                                         alignment=slide_content.get('title_alignment', PP_ALIGN.LEFT))
                    print("Updated title placeholder.")
                else:
                    # Fallback to adding a textbox at default position
                    left = Cm(2)
                    top = Cm(1)
                    width = Cm(29.867)
                    height = Cm(3)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    _set_text_properties(txBox.text_frame, title_text, title_font_props)
                    print(f"Added fallback title textbox at ({left.cm} cm, {top.cm} cm)")

            # Handle Body Text
            body_text = slide_content.get('text', '')
            body_font_props = slide_content.get('body_font', self.default_fonts['body'])
            body_position = slide_content.get('body_position', None)  # Optional

            if body_position:
                # Add body as a textbox with specified position
                left = Cm(body_position.get('left_cm', 2))
                top = Cm(body_position.get('top_cm', 4))
                width = Cm(body_position.get('width_cm', 29.867))
                height = Cm(body_position.get('height_cm', 13.05))
                txBox = slide.shapes.add_textbox(left, top, width, height)
                _set_text_properties(txBox.text_frame, body_text, body_font_props,
                                     alignment=body_position.get('alignment', PP_ALIGN.LEFT))
                print(f"Added custom body textbox at ({left.cm} cm, {top.cm} cm)")
            else:
                # Use existing body placeholder if available
                body_placeholder = _get_placeholder_by_type(slide, PP_PLACEHOLDER.OBJECT)
                if body_placeholder and body_placeholder.text_frame:
                    _set_text_properties(body_placeholder.text_frame, body_text, body_font_props,
                                         alignment=slide_content.get('body_alignment', PP_ALIGN.LEFT))
                    print("Updated body placeholder.")
                else:
                    # Fallback to adding a textbox at default position
                    left = Cm(2)
                    top = Cm(4)
                    width = Cm(29.867)
                    height = Cm(13.05)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    _set_text_properties(txBox.text_frame, body_text, body_font_props)
                    print(f"Added fallback body textbox at ({left.cm} cm, {top.cm} cm)")

            # Handle Columns (if any)
            if 'columns' in slide_content:
                columns = slide_content['columns']
                columns_font_props = slide_content.get('columns_font', self.default_fonts['body'])
                num_columns = len(columns)
                column_width = self.prs.slide_width / num_columns
                height = self.prs.slide_height - Cm(3.81)  # Adjust as needed

                for col_idx, column_text in enumerate(columns):
                    left = Cm(column_width.cm * col_idx)
                    top = Cm(3.81)
                    txBox = slide.shapes.add_textbox(left, top, column_width, height)
                    _set_text_properties(txBox.text_frame, column_text, columns_font_props)
                    print(f"Added column {col_idx + 1} textbox at ({left.cm} cm, {top.cm} cm)")

            # Handle Grid of Images
            if 'grid' in slide_content:
                grid = slide_content['grid']
                self._add_image_grid(slide, grid)
                print("Added image grid.")

            # Handle Table
            if 'table' in slide_content:
                # Add table from DataFrame
                table_content = slide_content['table']
                df = table_content['dataframe']
                position = table_content['position']
                left_mm = position.get('left_mm', 10)
                top_mm = position.get('top_mm', 50)
                width_mm = position.get('width_mm', 200)
                height_mm = position.get('height_mm', 100)
                header_font_props = table_content.get('header_font', self.default_fonts['header'])
                cell_font_props = table_content.get('cell_font', self.default_fonts['cell'])
                _add_table(slide, df, left_mm, top_mm, width_mm, height_mm, header_font_props, cell_font_props)
                print("Added table.")

            # Handle Additional Content (e.g., Images, Shapes)
            # Add any other content handling here as needed.

        try:
            self.prs.save(output_file)
            print(f"Presentation saved successfully to '{output_file}'.")
        except Exception as e:
            print(f"Failed to save presentation: {e}")

    @staticmethod
    def _add_image_grid(slide, grid):
        """
        Adds a grid of images to the given slide within a specified area.

        Parameters:
            slide (Slide): The slide to add images to.
            grid (dict): Configuration for the image grid, should contain:
                - image_configs (list): List of dicts with image paths as keys and (width_cm, height_cm) as values.
                - columns (int): Number of columns in the grid.
                - rows (int): Number of rows in the grid.
                - position_cm (dict): Dictionary with 'left_cm' and 'top_cm' specifying the grid's top-left position.
                - size_cm (dict): Dictionary with 'width_cm' and 'height_cm' specifying the grid's size.
        """
        # Extract image configurations
        image_configs = grid.get('image_configs', [])
        num_images = len(image_configs)

        if num_images == 0:
            print("No images to add to grid.")
            return

        # Determine the number of columns and rows
        cols = grid.get('columns', max(1, int(num_images ** 0.5)))
        rows = grid.get('rows', (num_images + cols - 1) // cols)

        # Get grid position and size
        position = grid.get('position_cm', {'left_cm': 1.0, 'top_cm': 1.0})
        size = grid.get('size_cm', {'width_cm': 20.0, 'height_cm': 15.0})

        left = Cm(position.get('left_cm', 1.0))
        top = Cm(position.get('top_cm', 1.0))
        grid_width = Cm(size.get('width_cm', 220.0))
        grid_height = Cm(size.get('height_cm', 156.0))

        # Calculate cell size so that images perfectly fit without spaces
        # Ensure cell_width and cell_height are Length objects
        cell_width = Cm(grid_width.cm / cols)
        cell_height = Cm(grid_height.cm / rows)

        for idx, img_config in enumerate(image_configs):
            if idx >= rows * cols:
                print(f"Skipping image {idx + 1}: exceeds grid capacity of {rows * cols} images.")
                break

            try:
                img_path, (width_cm, height_cm) = list(img_config.items())[0]
            except (IndexError, ValueError) as e:
                print(f"Invalid image configuration at index {idx}: {e}")
                continue

            try:
                # Convert cm to Length objects
                img_width = Cm(width_cm)
                img_height = Cm(height_cm)

                # Calculate scaling factor to fit within cell
                scale = min(cell_width / img_width, cell_height / img_height)
                img_width_scaled = Cm(width_cm * scale)
                img_height_scaled = Cm(height_cm * scale)

                # Calculate image position within the grid cell
                col = idx % cols
                row = idx // cols

                # Position within the grid cell
                # Ensure that col * cell_width.cm results in a float, then wrap back into Cm
                cell_left = left + Cm(col * cell_width.cm)
                cell_top = top + Cm(row * cell_height.cm)

                # Center the image within the cell
                # Calculate the delta in cm, then wrap into Cm
                delta_left = Cm((cell_width.cm - img_width_scaled.cm) / 2)
                delta_top = Cm((cell_height.cm - img_height_scaled.cm) / 2)
                img_left = cell_left + delta_left
                img_top = cell_top + delta_top

                # Add picture without any padding
                slide.shapes.add_picture(
                    img_path,
                    img_left,
                    img_top,
                    width=img_width_scaled,
                    height=img_height_scaled
                )
            except FileNotFoundError:
                logging.error(f"Image file not found: '{img_path}'")
            except Exception as e:
                logging.error(f"Error adding image '{img_path}': {e}")
