from pptx.enum.text import PP_ALIGN
from pptx.util import Mm, Pt


def _set_text_properties(text_frame, text, font_props):
    p = text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    font = p.runs[0].font
    font.name = font_props.get('name', 'Arial')
    font.size = font_props.get('size', Pt(14))
    font.bold = font_props.get('bold', False)


def _get_placeholder_by_idx(slide, idx):
    for placeholder in slide.placeholders:
        if placeholder.placeholder_format.idx == idx:
            return placeholder
    return None


def _add_table(slide, df, left_mm, top_mm, width_mm, height_mm, header_font_props, cell_font_props):
    left = Mm(left_mm)
    top = Mm(top_mm)
    width = Mm(width_mm)
    height = Mm(height_mm)

    rows = df.shape[0] + 1  # Add one for the header row
    cols = df.shape[1]

    # Add table shape
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths (equal widths for simplicity)
    for col_idx in range(cols):
        table.columns[col_idx].width = int(width / cols)

    # Write the header row
    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                font = run.font
                font.name = header_font_props.get('name', 'ＭＳ Ｐゴシック (本文)')
                font.size = header_font_props.get('size', Pt(11))
                font.bold = header_font_props.get('bold', True)

    # Write the data rows
    for row_idx, (index, row) in enumerate(df.iterrows(), start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    font = run.font
                    font.name = cell_font_props.get('name', 'ＭＳ Ｐゴシック (本文)')
                    font.size = cell_font_props.get('size', Pt(11))
                    font.bold = cell_font_props.get('bold', False)
