from math import floor

from pptx import Presentation
from pptx.util import Inches, Pt, Mm, Cm
from pptx.enum.text import PP_ALIGN

from .utils import _set_text_properties, _get_placeholder_by_idx, _add_table


class CustomPresentationBuilder:
    def __init__(self, slides_content_list):
        self.slides_content_list = slides_content_list
        self.prs = Presentation()
        # Set slide size to 16:9 aspect ratio in centimeters
        self.prs.slide_width = Cm(33.867)  # 13.3333 inches * 2.54 cm/in ≈ 33.867 cm
        self.prs.slide_height = Cm(19.05)  # 7.5 inches * 2.54 cm/in = 19.05 cm
        self.layout_mapping = self._create_layout_mapping()
        self.default_fonts = {
            'title': {'name': 'Arial', 'size': Pt(26.7), 'bold': False},
            'body': {'name': 'Arial', 'size': Pt(14), 'bold': False},
            'header': {'name': 'ＭＳ Ｐゴシック (本文)', 'size': Pt(11), 'bold': True},
            'cell': {'name': 'ＭＳ Ｐゴシック (本文)', 'size': Pt(11), 'bold': False},
        }

    def _create_layout_mapping(self):
        layout_mapping = {}
        for layout in self.prs.slide_layouts:
            name = layout.name.upper().replace(' ', '_')
            layout_mapping[name] = layout
        return layout_mapping

    def create_presentation(self, output_file):
        for slide_content in self.slides_content_list:
            layout_name = slide_content.get('layout', 'TITLE_AND_CONTENT').upper()
            layout = self.layout_mapping.get(layout_name, self.prs.slide_layouts[1])
            slide = self.prs.slides.add_slide(layout)

            # Set title
            title_placeholder = slide.shapes.title
            title_font_props = slide_content.get('title_font', self.default_fonts['title'])
            _set_text_properties(title_placeholder.text_frame, slide_content.get('title', ''), title_font_props)

            # Handle different layouts
            if 'columns' in slide_content:
                columns = slide_content['columns']
                body_font_props = slide_content.get('body_font', self.default_fonts['body'])
                if layout_name == 'TWO_CONTENT' and len(columns) == 2:
                    left_placeholder = _get_placeholder_by_idx(slide, 1)
                    right_placeholder = _get_placeholder_by_idx(slide, 2)
                    if left_placeholder and right_placeholder:
                        _set_text_properties(left_placeholder.text_frame, columns[0], body_font_props)
                        _set_text_properties(right_placeholder.text_frame, columns[1], body_font_props)
                    else:
                        print("Error: Expected placeholders not found in slide layout.")
                else:
                    # Manually create columns
                    shapes = slide.shapes
                    num_columns = len(columns)
                    column_width = self.prs.slide_width / num_columns
                    height = self.prs.slide_height - Cm(3.81)  # 1.5 inches * 2.54 cm/in ≈ 3.81 cm
                    for idx, column_text in enumerate(columns):
                        left = column_width * idx
                        top = Cm(3.81)  # 1.5 inches ≈ 3.81 cm
                        txBox = shapes.add_textbox(left, top, column_width, height)
                        _set_text_properties(txBox.text_frame, column_text, body_font_props)

            if 'text' in slide_content:
                # Single body text
                body_placeholder = _get_placeholder_by_idx(slide, 1)
                body_font_props = slide_content.get('body_font', self.default_fonts['body'])
                if body_placeholder:
                    _set_text_properties(body_placeholder.text_frame, slide_content['text'], body_font_props)
                else:
                    # Create a new textbox
                    shapes = slide.shapes
                    left = Cm(1.34)  # 1 inch ≈ 2.54 cm
                    top = Cm(0.89)  # 2 inches ≈ 5.08 cm
                    width = self.prs.slide_width - Cm(1.34)  # Subtracting 2 inches on both sides
                    height = self.prs.slide_height - Cm(0.89)  # Subtracting 3 inches from height
                    txBox = shapes.add_textbox(left, top, width, height)
                    _set_text_properties(txBox.text_frame, slide_content['text'], body_font_props)

            if 'grid' in slide_content:
                # Grid of images
                grid = slide_content['grid']
                self._add_image_grid(slide, grid)

            if 'table' in slide_content:
                # Add table from DataFrame
                table_content = slide_content['table']
                df = table_content['dataframe']
                position = table_content['position']
                left_mm = position.get('left_mm', 10)
                top_mm = position.get('top_mm', 50)
                width_mm = position.get('width_mm', 200)
                height_mm = position.get('height_mm', 100)
                header_font_props = table_content.get('header_font', self.default_fonts['header'])
                cell_font_props = table_content.get('cell_font', self.default_fonts['cell'])
                _add_table(slide, df, left_mm, top_mm, width_mm, height_mm, header_font_props, cell_font_props)
            else:
                # No additional content
                pass
        try:
            self.prs.save(output_file)
            print(f"Presentation saved successfully to '{output_file}'.")
        except Exception as e:
            print(f"Failed to save presentation: {e}")

    @staticmethod
    def _add_image_grid(slide, grid):
        """
        Adds a grid of images to the given slide within a specified area.

        Parameters:
            slide (Slide): The slide to add images to.
            grid (dict): Configuration for the image grid, should contain:
                - image_configs (list): List of dicts with image paths as keys and (width_cm, height_cm) as values.
                - columns (int): Number of columns in the grid.
                - rows (int): Number of rows in the grid.
                - position_cm (dict): Dictionary with 'left_cm' and 'top_cm' specifying the grid's top-left position.
                - size_cm (dict): Dictionary with 'width_cm' and 'height_cm' specifying the grid's size.
        """
        # Extract image configurations
        image_configs = grid.get('image_configs', [])
        num_images = len(image_configs)

        if num_images == 0:
            print("No images to add to grid.")
            return

        # Determine the number of columns and rows
        cols = grid.get('columns', max(1, int(num_images ** 0.5)))
        rows = grid.get('rows', (num_images + cols - 1) // cols)

        # Get grid position and size
        position = grid.get('position_cm', {'left_cm': 1.0, 'top_cm': 1.0})
        size = grid.get('size_cm', {'width_cm': 20.0, 'height_cm': 15.0})

        left = Cm(position.get('left_cm', 1.0))
        top = Cm(position.get('top_cm', 1.0))
        grid_width = Cm(size.get('width_cm', 220.0))
        grid_height = Cm(size.get('height_cm', 156.0))

        # Calculate cell size so that images perfectly fit without spaces
        # Ensure cell_width and cell_height are Length objects
        cell_width = Cm(grid_width.cm / cols)
        cell_height = Cm(grid_height.cm / rows)

        for idx, img_config in enumerate(image_configs):
            if idx >= rows * cols:
                print(f"Skipping image {idx + 1}: exceeds grid capacity of {rows * cols} images.")
                break

            try:
                img_path, (width_cm, height_cm) = list(img_config.items())[0]
            except (IndexError, ValueError) as e:
                print(f"Invalid image configuration at index {idx}: {e}")
                continue

            try:
                # Convert cm to Length objects
                img_width = Cm(width_cm)
                img_height = Cm(height_cm)

                # Calculate scaling factor to fit within cell
                scale = min(cell_width / img_width, cell_height / img_height)
                img_width_scaled = Cm(width_cm * scale)
                img_height_scaled = Cm(height_cm * scale)

                # Calculate image position within the grid cell
                col = idx % cols
                row = idx // cols

                # Position within the grid cell
                # Ensure that col * cell_width.cm results in a float, then wrap back into Cm
                cell_left = left + Cm(col * cell_width.cm)
                cell_top = top + Cm(row * cell_height.cm)

                # Center the image within the cell
                # Calculate the delta in cm, then wrap into Cm
                delta_left = Cm((cell_width.cm - img_width_scaled.cm) / 2)
                delta_top = Cm((cell_height.cm - img_height_scaled.cm) / 2)
                img_left = cell_left + delta_left
                img_top = cell_top + delta_top

                # Add picture without any padding
                slide.shapes.add_picture(
                    img_path,
                    img_left,
                    img_top,
                    width=img_width_scaled,
                    height=img_height_scaled
                )
            except FileNotFoundError:
                print(f"Image file not found: '{img_path}'")
            except Exception as e:
                print(f"Error adding image '{img_path}': {e}")
